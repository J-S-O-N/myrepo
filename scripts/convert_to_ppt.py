#!/usr/bin/env python3
"""
Convert Markdown presentation to PowerPoint
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_markdown(md_file):
    """Parse markdown file and extract slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide separator (---) with any whitespace
    slides_content = re.split(r'\n\s*---\s*\n', content)

    slides = []
    for slide_content in slides_content:
        slide_content = slide_content.strip()
        # Skip empty slides and header
        if not slide_content or slide_content.startswith('# BankApp - Technical Presentation'):
            continue

        slides.append(slide_content)

    return slides

def parse_slide_content(content):
    """Parse individual slide content"""
    lines = content.split('\n')
    title = ""
    body_lines = []
    code_blocks = []

    in_code_block = False
    code_lines = []

    for line in lines:
        line = line.strip()

        # Handle code blocks
        if line.startswith('```'):
            if in_code_block:
                code_blocks.append('\n'.join(code_lines))
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        # Extract title (## or ###)
        if line.startswith('## '):
            title = line[3:].strip()
        elif line.startswith('### '):
            if not title:
                title = line[4:].strip()
            else:
                body_lines.append(line)
        elif line and not line.startswith('#'):
            body_lines.append(line)

    return title, body_lines, code_blocks

def clean_markdown(text):
    """Remove markdown formatting"""
    # Remove bold
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remove italic
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Remove inline code
    text = re.sub(r'`(.*?)`', r'\1', text)
    # Remove links but keep text
    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
    return text

def create_presentation(slides_data, output_file):
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define colors
    TITLE_COLOR = RGBColor(2, 132, 199)  # Blue
    TEXT_COLOR = RGBColor(55, 65, 81)    # Dark gray
    ACCENT_COLOR = RGBColor(16, 185, 129) # Green

    for slide_content in slides_data:
        title, body_lines, code_blocks = parse_slide_content(slide_content)

        if not title:
            continue

        # Add slide
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = clean_markdown(title)
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

        # Add body content
        if len(slide.shapes) > 1:
            body_shape = slide.shapes[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            # Add body lines
            for i, line in enumerate(body_lines[:15]):  # Limit to 15 lines per slide
                if not line.strip():
                    continue

                cleaned_line = clean_markdown(line)

                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                # Check if it's a bullet point
                if line.startswith('- ') or line.startswith('* ') or line.startswith('✅') or line.startswith('⏳'):
                    p.text = cleaned_line.lstrip('- *')
                    p.level = 0
                elif line.startswith('  - ') or line.startswith('  * '):
                    p.text = cleaned_line.lstrip('- *')
                    p.level = 1
                else:
                    p.text = cleaned_line
                    p.level = 0

                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR

            # Add code blocks if any
            for code_block in code_blocks[:2]:  # Limit to 2 code blocks per slide
                p = text_frame.add_paragraph()
                p.text = code_block[:500]  # Limit code length
                p.font.size = Pt(12)
                p.font.name = 'Courier New'
                p.font.color.rgb = RGBColor(88, 110, 117)
                p.level = 0

    # Save presentation
    prs.save(output_file)
    print(f"✅ PowerPoint presentation created: {output_file}")

if __name__ == '__main__':
    import sys

    md_file = 'docs/BankApp_Presentation.md'
    output_file = 'docs/BankApp_Presentation.pptx'

    print(f"📄 Reading: {md_file}")
    slides = parse_markdown(md_file)
    print(f"📊 Found {len(slides)} slides")

    print(f"🎨 Creating PowerPoint presentation...")
    create_presentation(slides, output_file)
    print(f"✅ Done! Presentation saved to: {output_file}")
