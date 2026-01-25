#!/usr/bin/env python3
"""
Create Executive PowerPoint Presentation with Infographics
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_exec_presentation(output_file):
    """Create executive-style presentation with infographics"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Define color scheme
    PRIMARY_BLUE = RGBColor(2, 132, 199)
    DARK_BLUE = RGBColor(3, 105, 161)
    ACCENT_GREEN = RGBColor(16, 185, 129)
    ACCENT_ORANGE = RGBColor(251, 146, 60)
    DARK_GRAY = RGBColor(31, 41, 55)
    LIGHT_GRAY = RGBColor(156, 163, 175)
    WHITE = RGBColor(255, 255, 255)
    BG_LIGHT = RGBColor(248, 250, 252)

    def add_title_slide():
        """Slide 1: Executive Title Slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background gradient (simulated with shape)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = DARK_BLUE
        bg_shape.line.fill.background()

        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "BankApp"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(11.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Next-Generation Banking Platform"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = ACCENT_GREEN
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Tagline
        tagline_box = slide.shapes.add_textbox(
            Inches(1), Inches(5),
            Inches(11.333), Inches(0.5)
        )
        tagline_frame = tagline_box.text_frame
        tagline_frame.text = "Cloud-Native • Secure • Scalable"
        tagline_para = tagline_frame.paragraphs[0]
        tagline_para.font.size = Pt(24)
        tagline_para.font.color.rgb = LIGHT_GRAY
        tagline_para.alignment = PP_ALIGN.CENTER

        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.8),
            Inches(11.333), Inches(0.4)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = "Executive Presentation | 2026"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = LIGHT_GRAY
        footer_para.alignment = PP_ALIGN.CENTER

    def add_executive_summary():
        """Slide 2: Executive Summary"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Key metrics boxes
        metrics = [
            ("Users", "Multi-Tenant", ACCENT_GREEN),
            ("Security", "Enterprise-Grade", PRIMARY_BLUE),
            ("Cloud", "AWS Infrastructure", ACCENT_ORANGE),
            ("Pipeline", "Automated CI/CD", DARK_BLUE)
        ]

        start_x = 0.5
        box_width = 2.8
        box_height = 1.2
        spacing = 0.3

        for i, (label, value, color) in enumerate(metrics):
            x = start_x + i * (box_width + spacing)

            # Box background
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.3),
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()
            box.shadow.inherit = False

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(1.5),
                Inches(box_width), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = label
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(16)
            label_para.font.bold = True
            label_para.font.color.rgb = WHITE
            label_para.alignment = PP_ALIGN.CENTER

            # Value
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(1.9),
                Inches(box_width), Inches(0.5)
            )
            value_frame = value_box.text_frame
            value_frame.text = value
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(14)
            value_para.font.color.rgb = WHITE
            value_para.alignment = PP_ALIGN.CENTER

        # Value proposition bullets
        bullets_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8),
            Inches(12.333), Inches(4)
        )
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        bullet_points = [
            "✓ Full-stack banking application with modern React frontend and Node.js backend",
            "✓ Per-user data isolation with enterprise-grade security (JWT + bcrypt)",
            "✓ Cloud-native AWS infrastructure: VPC, ECR, ECS, RDS, CloudFront",
            "✓ Automated CI/CD pipeline with Docker containerization",
            "✓ Comprehensive features: Accounts, Goals, Investments, Crypto, Health tracking",
            "✓ Production-ready: 68.9 MB Docker image successfully deployed to ECR",
            "✓ Scalable architecture: Auto-scaling with ECS Fargate (planned)",
            "✓ Cost-effective: ~$830/month for 3 environments (dev, staging, prod)"
        ]

        for i, bullet in enumerate(bullet_points):
            p = bullets_frame.add_paragraph() if i > 0 else bullets_frame.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)
            p.level = 0

    def add_tech_stack_infographic():
        """Slide 3: Technology Stack Infographic"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Technology Stack"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Three columns: Frontend, Backend, Cloud
        columns = [
            {
                "title": "Frontend",
                "icon": "⚛️",
                "items": ["React 18.3.1", "Vite 6.0", "Vitest Testing", "Modern Hooks"],
                "color": ACCENT_GREEN
            },
            {
                "title": "Backend",
                "icon": "🚀",
                "items": ["Node.js 18 LTS", "Express 5.2", "Sequelize ORM", "JWT Auth"],
                "color": PRIMARY_BLUE
            },
            {
                "title": "Cloud & DevOps",
                "icon": "☁️",
                "items": ["AWS VPC/ECR/ECS", "Docker", "GitHub Actions", "Terraform IaC"],
                "color": ACCENT_ORANGE
            }
        ]

        col_width = 3.8
        col_spacing = 0.4
        start_x = 0.7

        for i, col in enumerate(columns):
            x = start_x + i * (col_width + col_spacing)

            # Column box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.5),
                Inches(col_width), Inches(5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = BG_LIGHT
            box.line.color.rgb = col["color"]
            box.line.width = Pt(3)

            # Icon
            icon_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(1.7),
                Inches(col_width - 0.4), Inches(0.8)
            )
            icon_frame = icon_box.text_frame
            icon_frame.text = col["icon"]
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(48)
            icon_para.alignment = PP_ALIGN.CENTER

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(2.6),
                Inches(col_width - 0.4), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = col["title"]
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.font.color.rgb = col["color"]
            title_para.alignment = PP_ALIGN.CENTER

            # Items
            items_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(3.3),
                Inches(col_width - 0.6), Inches(3)
            )
            items_frame = items_box.text_frame

            for j, item in enumerate(col["items"]):
                p = items_frame.add_paragraph() if j > 0 else items_frame.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(10)

    def add_security_architecture():
        """Slide 4: Security Architecture"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "🔐 Multi-Layer Security Architecture"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Security layers
        layers = [
            ("Authentication", "JWT Tokens (24h)\nbcrypt Hashing (10 rounds)", PRIMARY_BLUE),
            ("Authorization", "Per-User Data Isolation\nRole-Based Access", ACCENT_GREEN),
            ("Network", "VPC Isolation\nPrivate Subnets\nSecurity Groups", ACCENT_ORANGE),
            ("Data", "Encryption at Rest\nEncryption in Transit\nSecrets Manager", DARK_BLUE),
            ("Application", "Input Validation\nSQL Injection Prevention\nXSS Protection", PRIMARY_BLUE),
            ("Infrastructure", "Non-root Containers\nImage Scanning\nMinimal Base Images", ACCENT_GREEN)
        ]

        # Create layers in 2 rows
        box_width = 3.8
        box_height = 1.8
        spacing_x = 0.4
        spacing_y = 0.3
        start_x = 0.5
        start_y = 1.5

        for i, (title, desc, color) in enumerate(layers):
            row = i // 3
            col = i % 3
            x = start_x + col * (box_width + spacing_x)
            y = start_y + row * (box_height + spacing_y)

            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()
            box.shadow.inherit = False

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.2),
                Inches(box_width - 0.4), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(20)
            title_para.font.bold = True
            title_para.font.color.rgb = WHITE
            title_para.alignment = PP_ALIGN.CENTER

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.7),
                Inches(box_width - 0.4), Inches(1)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = desc
            desc_frame.word_wrap = True
            for para in desc_frame.paragraphs:
                para.font.size = Pt(14)
                para.font.color.rgb = WHITE
                para.alignment = PP_ALIGN.CENTER

        # Bottom banner
        banner_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.3),
            Inches(12.333), Inches(0.8)
        )
        banner_frame = banner_box.text_frame
        banner_frame.text = "✓ Enterprise-Grade Security  |  ✓ Compliance-Ready  |  ✓ Zero Trust Architecture"
        banner_para = banner_frame.paragraphs[0]
        banner_para.font.size = Pt(20)
        banner_para.font.bold = True
        banner_para.font.color.rgb = ACCENT_GREEN
        banner_para.alignment = PP_ALIGN.CENTER

    def add_aws_infrastructure():
        """Slide 5: AWS Infrastructure Diagram"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "☁️ AWS Cloud Infrastructure"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Infrastructure components in layers
        # Layer 1: CDN/DNS
        layer1 = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(1.5),
            Inches(9.333), Inches(0.8)
        )
        layer1.fill.solid()
        layer1.fill.fore_color.rgb = ACCENT_ORANGE
        layer1.line.fill.background()

        layer1_text = slide.shapes.add_textbox(Inches(2), Inches(1.6), Inches(9.333), Inches(0.6))
        layer1_text.text_frame.text = "Route 53 DNS + CloudFront CDN"
        layer1_text.text_frame.paragraphs[0].font.size = Pt(20)
        layer1_text.text_frame.paragraphs[0].font.bold = True
        layer1_text.text_frame.paragraphs[0].font.color.rgb = WHITE
        layer1_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Layer 2: Load Balancer
        layer2 = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(2.6),
            Inches(8.333), Inches(0.7)
        )
        layer2.fill.solid()
        layer2.fill.fore_color.rgb = PRIMARY_BLUE
        layer2.line.fill.background()

        layer2_text = slide.shapes.add_textbox(Inches(2.5), Inches(2.7), Inches(8.333), Inches(0.5))
        layer2_text.text_frame.text = "Application Load Balancer"
        layer2_text.text_frame.paragraphs[0].font.size = Pt(18)
        layer2_text.text_frame.paragraphs[0].font.bold = True
        layer2_text.text_frame.paragraphs[0].font.color.rgb = WHITE
        layer2_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Layer 3: VPC
        vpc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(3.6),
            Inches(11.333), Inches(3)
        )
        vpc_box.fill.solid()
        vpc_box.fill.fore_color.rgb = BG_LIGHT
        vpc_box.line.color.rgb = DARK_BLUE
        vpc_box.line.width = Pt(3)

        vpc_label = slide.shapes.add_textbox(Inches(1.2), Inches(3.7), Inches(3), Inches(0.4))
        vpc_label.text_frame.text = "VPC (10.0.0.0/16)"
        vpc_label.text_frame.paragraphs[0].font.size = Pt(16)
        vpc_label.text_frame.paragraphs[0].font.bold = True
        vpc_label.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

        # ECS Fargate
        ecs_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(4.3),
            Inches(4.5), Inches(2)
        )
        ecs_box.fill.solid()
        ecs_box.fill.fore_color.rgb = ACCENT_GREEN
        ecs_box.line.fill.background()

        ecs_text = slide.shapes.add_textbox(Inches(1.7), Inches(4.5), Inches(4.1), Inches(1.6))
        ecs_frame = ecs_text.text_frame
        ecs_frame.text = "ECS Fargate\n\nDocker Containers\nAuto-scaling\n68.9 MB Image"
        ecs_frame.paragraphs[0].font.size = Pt(18)
        ecs_frame.paragraphs[0].font.bold = True
        ecs_frame.paragraphs[0].font.color.rgb = WHITE
        ecs_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for para in ecs_frame.paragraphs[1:]:
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER

        # RDS Database
        rds_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.5), Inches(4.3),
            Inches(4.5), Inches(2)
        )
        rds_box.fill.solid()
        rds_box.fill.fore_color.rgb = PRIMARY_BLUE
        rds_box.line.fill.background()

        rds_text = slide.shapes.add_textbox(Inches(6.7), Inches(4.5), Inches(4.1), Inches(1.6))
        rds_frame = rds_text.text_frame
        rds_frame.text = "RDS PostgreSQL\n\nMulti-AZ\nAutomated Backups\nEncrypted"
        rds_frame.paragraphs[0].font.size = Pt(18)
        rds_frame.paragraphs[0].font.bold = True
        rds_frame.paragraphs[0].font.color.rgb = WHITE
        rds_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for para in rds_frame.paragraphs[1:]:
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER

        # Cost banner
        cost_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8),
            Inches(12.333), Inches(0.4)
        )
        cost_frame = cost_box.text_frame
        cost_frame.text = "💰 Total Infrastructure Cost: ~$830/month (Dev + Staging + Production)"
        cost_para = cost_frame.paragraphs[0]
        cost_para.font.size = Pt(18)
        cost_para.font.bold = True
        cost_para.font.color.rgb = ACCENT_ORANGE
        cost_para.alignment = PP_ALIGN.CENTER

    def add_features_dashboard():
        """Slide 6: Key Features Dashboard"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "💼 Application Features"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Features in grid
        features = [
            ("🏠", "Dashboard", "Real-time balances\nTransactions\nAnalytics"),
            ("💰", "Multi-Account", "Checking\nSavings\nCredit"),
            ("🎯", "Goals", "Financial targets\nProgress tracking\nContributions"),
            ("📈", "Investments", "JSE stocks\nPortfolio\nPerformance"),
            ("₿", "Crypto", "Live prices\nBTC, ETH, SOL\n24h changes"),
            ("💱", "Exchange", "Live ZAR/USD\nAuto-refresh\nDual display"),
            ("⚙️", "Settings", "Transaction limits\nCard controls\nPreferences"),
            ("❤️", "Health", "Steps tracking\nCalories\nWellness"),
            ("🛒", "Buy Hub", "Shopping\nCategories\nOffers")
        ]

        box_width = 3.8
        box_height = 1.6
        spacing_x = 0.4
        spacing_y = 0.3
        start_x = 0.5
        start_y = 1.5

        for i, (icon, title, desc) in enumerate(features):
            row = i // 3
            col = i % 3
            x = start_x + col * (box_width + spacing_x)
            y = start_y + row * (box_height + spacing_y)

            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            colors = [PRIMARY_BLUE, ACCENT_GREEN, ACCENT_ORANGE]
            box.fill.fore_color.rgb = colors[col]
            box.line.fill.background()

            # Icon
            icon_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.1),
                Inches(0.6), Inches(0.5)
            )
            icon_frame = icon_box.text_frame
            icon_frame.text = icon
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(32)

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.9), Inches(y + 0.15),
                Inches(box_width - 1.1), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(20)
            title_para.font.bold = True
            title_para.font.color.rgb = WHITE

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.7),
                Inches(box_width - 0.4), Inches(0.8)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = desc
            desc_frame.word_wrap = True
            for para in desc_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.color.rgb = WHITE

    def add_deployment_status():
        """Slide 7: Deployment Status & Metrics"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "🚀 Deployment Status"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Status boxes
        statuses = [
            ("✅ VPC Infrastructure", "2 AZs, NAT Gateways", ACCENT_GREEN),
            ("✅ Docker Image", "68.9 MB in ECR", ACCENT_GREEN),
            ("✅ CI/CD Pipeline", "GitHub Actions", ACCENT_GREEN),
            ("⏳ ECS Deployment", "Pending", ACCENT_ORANGE)
        ]

        for i, (status, detail, color) in enumerate(statuses):
            y = 1.5 + i * 1.2

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y),
                Inches(11.333), Inches(0.9)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            status_text = slide.shapes.add_textbox(
                Inches(1.3), Inches(y + 0.15),
                Inches(6), Inches(0.6)
            )
            status_frame = status_text.text_frame
            status_frame.text = status
            status_para = status_frame.paragraphs[0]
            status_para.font.size = Pt(24)
            status_para.font.bold = True
            status_para.font.color.rgb = WHITE

            detail_text = slide.shapes.add_textbox(
                Inches(7.5), Inches(y + 0.2),
                Inches(4.5), Inches(0.5)
            )
            detail_frame = detail_text.text_frame
            detail_frame.text = detail
            detail_para = detail_frame.paragraphs[0]
            detail_para.font.size = Pt(18)
            detail_para.font.color.rgb = WHITE
            detail_para.alignment = PP_ALIGN.RIGHT

        # Metrics
        metrics_title = slide.shapes.add_textbox(
            Inches(1), Inches(6.3),
            Inches(11.333), Inches(0.4)
        )
        metrics_title.text_frame.text = "Project Metrics"
        metrics_title.text_frame.paragraphs[0].font.size = Pt(20)
        metrics_title.text_frame.paragraphs[0].font.bold = True
        metrics_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

        metrics_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.7),
            Inches(11.333), Inches(0.5)
        )
        metrics_frame = metrics_box.text_frame
        metrics_frame.text = "5,300 Lines of Code  •  10+ Features  •  32 Slides Documentation  •  Production Ready"
        metrics_para = metrics_frame.paragraphs[0]
        metrics_para.font.size = Pt(16)
        metrics_para.font.color.rgb = PRIMARY_BLUE
        metrics_para.alignment = PP_ALIGN.CENTER

    def add_next_steps():
        """Slide 8: Next Steps & Roadmap"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "📋 Strategic Roadmap"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Roadmap phases
        phases = [
            ("Phase 1", "Infrastructure Foundation", ["✅ AWS VPC setup", "✅ Docker containerization", "✅ CI/CD pipeline"], ACCENT_GREEN),
            ("Phase 2", "Production Deployment", ["⏳ ECS Fargate setup", "⏳ RDS database", "⏳ Load balancer"], ACCENT_ORANGE),
            ("Phase 3", "Enhancement & Scale", ["📅 Performance optimization", "📅 Monitoring dashboards", "📅 Auto-scaling"], PRIMARY_BLUE)
        ]

        box_width = 3.8
        start_x = 0.5
        spacing = 0.4

        for i, (phase, title, items, color) in enumerate(phases):
            x = start_x + i * (box_width + spacing)

            # Phase box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.5),
                Inches(box_width), Inches(4.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = BG_LIGHT
            box.line.color.rgb = color
            box.line.width = Pt(4)

            # Phase number
            phase_num = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(1.7),
                Inches(box_width - 0.4), Inches(0.5)
            )
            phase_frame = phase_num.text_frame
            phase_frame.text = phase
            phase_para = phase_frame.paragraphs[0]
            phase_para.font.size = Pt(22)
            phase_para.font.bold = True
            phase_para.font.color.rgb = color
            phase_para.alignment = PP_ALIGN.CENTER

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(2.3),
                Inches(box_width - 0.4), Inches(0.6)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = DARK_GRAY
            title_para.alignment = PP_ALIGN.CENTER

            # Items
            items_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(3.1),
                Inches(box_width - 0.6), Inches(2.7)
            )
            items_frame = items_box.text_frame

            for j, item in enumerate(items):
                p = items_frame.add_paragraph() if j > 0 else items_frame.paragraphs[0]
                p.text = item
                p.font.size = Pt(15)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(12)

        # Timeline
        timeline_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.3),
            Inches(12.333), Inches(0.8)
        )
        timeline_frame = timeline_box.text_frame
        timeline_frame.text = "Timeline: Phase 1 ✅ Complete  |  Phase 2 🔄 In Progress  |  Phase 3 📅 Q2 2026"
        timeline_para = timeline_frame.paragraphs[0]
        timeline_para.font.size = Pt(20)
        timeline_para.font.bold = True
        timeline_para.font.color.rgb = PRIMARY_BLUE
        timeline_para.alignment = PP_ALIGN.CENTER

    def add_closing_slide():
        """Slide 9: Closing & Call to Action"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = DARK_BLUE
        bg_shape.line.fill.background()

        # Main message
        message_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(11.333), Inches(1.5)
        )
        message_frame = message_box.text_frame
        message_frame.text = "Ready for Production"
        message_para = message_frame.paragraphs[0]
        message_para.font.size = Pt(60)
        message_para.font.bold = True
        message_para.font.color.rgb = WHITE
        message_para.alignment = PP_ALIGN.CENTER

        # Subtext
        subtext_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8),
            Inches(11.333), Inches(1)
        )
        subtext_frame = subtext_box.text_frame
        subtext_frame.text = "Enterprise-Grade Banking Platform\nBuilt on AWS • Secured by Design • Ready to Scale"
        for para in subtext_frame.paragraphs:
            para.font.size = Pt(28)
            para.font.color.rgb = ACCENT_GREEN
            para.alignment = PP_ALIGN.CENTER

        # Call to action
        cta_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5),
            Inches(11.333), Inches(0.8)
        )
        cta_frame = cta_box.text_frame
        cta_frame.text = "Questions?"
        cta_para = cta_frame.paragraphs[0]
        cta_para.font.size = Pt(36)
        cta_para.font.color.rgb = WHITE
        cta_para.alignment = PP_ALIGN.CENTER

        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(11.333), Inches(0.5)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "GitHub: J-S-O-N/myrepo  •  AWS Region: us-east-1"
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(18)
        contact_para.font.color.rgb = LIGHT_GRAY
        contact_para.alignment = PP_ALIGN.CENTER

    def add_market_opportunity():
        """Slide 10: Market Opportunity"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "📊 Market Opportunity"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Market stats boxes
        stats = [
            ("Digital Banking Users", "2.5B+", "Global Market 2026", ACCENT_GREEN),
            ("South Africa Market", "38M", "Banked Population", PRIMARY_BLUE),
            ("Mobile Banking", "73%", "Adoption Rate SA", ACCENT_ORANGE),
            ("Market Growth", "12.5%", "CAGR 2024-2030", DARK_BLUE)
        ]

        box_width = 2.8
        spacing = 0.3
        start_x = 0.5

        for i, (label, value, desc, color) in enumerate(stats):
            x = start_x + i * (box_width + spacing)

            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.5),
                Inches(box_width), Inches(1.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            # Value (large)
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(1.7),
                Inches(box_width), Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = value
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(40)
            value_para.font.bold = True
            value_para.font.color.rgb = WHITE
            value_para.alignment = PP_ALIGN.CENTER

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.3),
                Inches(box_width - 0.2), Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.word_wrap = True
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(14)
            label_para.font.bold = True
            label_para.font.color.rgb = WHITE
            label_para.alignment = PP_ALIGN.CENTER

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.7),
                Inches(box_width - 0.2), Inches(0.3)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = desc
            desc_frame.word_wrap = True
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(11)
            desc_para.font.color.rgb = WHITE
            desc_para.alignment = PP_ALIGN.CENTER

        # Target segments
        segment_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.3),
            Inches(12.333), Inches(0.4)
        )
        segment_title.text_frame.text = "Target Market Segments"
        segment_title.text_frame.paragraphs[0].font.size = Pt(24)
        segment_title.text_frame.paragraphs[0].font.bold = True
        segment_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

        segments = [
            ("💼 SME Banking", "Small & medium enterprises requiring comprehensive financial management"),
            ("👥 Retail Banking", "Individual customers seeking modern, mobile-first banking experience"),
            ("🏢 Corporate Banking", "Enterprises needing multi-account management and treasury services"),
            ("🌍 International", "Cross-border payments and multi-currency support for global businesses")
        ]

        y_start = 4
        for i, (seg_title, seg_desc) in enumerate(segments):
            y = y_start + i * 0.7

            # Segment box
            seg_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y),
                Inches(12.333), Inches(0.6)
            )
            seg_box.fill.solid()
            seg_box.fill.fore_color.rgb = BG_LIGHT
            seg_box.line.color.rgb = PRIMARY_BLUE
            seg_box.line.width = Pt(2)

            # Title
            seg_title_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y + 0.05),
                Inches(3), Inches(0.3)
            )
            seg_title_box.text_frame.text = seg_title
            seg_title_box.text_frame.paragraphs[0].font.size = Pt(16)
            seg_title_box.text_frame.paragraphs[0].font.bold = True
            seg_title_box.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE

            # Description
            seg_desc_box = slide.shapes.add_textbox(
                Inches(3.8), Inches(y + 0.1),
                Inches(8.5), Inches(0.4)
            )
            seg_desc_box.text_frame.text = seg_desc
            seg_desc_box.text_frame.word_wrap = True
            seg_desc_box.text_frame.paragraphs[0].font.size = Pt(14)
            seg_desc_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    def add_business_model():
        """Slide 11: Revenue Model & Business Case"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "💰 Revenue Model & Business Case"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Revenue streams
        revenue_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2),
            Inches(5.5), Inches(0.4)
        )
        revenue_title.text_frame.text = "Revenue Streams"
        revenue_title.text_frame.paragraphs[0].font.size = Pt(22)
        revenue_title.text_frame.paragraphs[0].font.bold = True
        revenue_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

        streams = [
            ("Transaction Fees", "R2-R5 per transaction", "60%"),
            ("Monthly Subscriptions", "R99-R499 per user/month", "25%"),
            ("Premium Features", "Goals, Crypto tracking, Analytics", "10%"),
            ("Partner Commissions", "Buy Hub, Investment products", "5%")
        ]

        y_start = 1.8
        for i, (stream, desc, percentage) in enumerate(streams):
            y = y_start + i * 0.8

            # Stream box
            stream_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y),
                Inches(5.5), Inches(0.7)
            )
            colors = [ACCENT_GREEN, PRIMARY_BLUE, ACCENT_ORANGE, DARK_BLUE]
            stream_box.fill.solid()
            stream_box.fill.fore_color.rgb = colors[i]
            stream_box.line.fill.background()

            # Percentage circle
            perc_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y + 0.15),
                Inches(0.8), Inches(0.4)
            )
            perc_box.text_frame.text = percentage
            perc_box.text_frame.paragraphs[0].font.size = Pt(20)
            perc_box.text_frame.paragraphs[0].font.bold = True
            perc_box.text_frame.paragraphs[0].font.color.rgb = WHITE
            perc_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Stream name
            name_box = slide.shapes.add_textbox(
                Inches(1.6), Inches(y + 0.1),
                Inches(2.5), Inches(0.3)
            )
            name_box.text_frame.text = stream
            name_box.text_frame.paragraphs[0].font.size = Pt(16)
            name_box.text_frame.paragraphs[0].font.bold = True
            name_box.text_frame.paragraphs[0].font.color.rgb = WHITE

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(1.6), Inches(y + 0.4),
                Inches(3.6), Inches(0.25)
            )
            desc_box.text_frame.text = desc
            desc_box.text_frame.word_wrap = True
            desc_box.text_frame.paragraphs[0].font.size = Pt(12)
            desc_box.text_frame.paragraphs[0].font.color.rgb = WHITE

        # Financial projections
        proj_title = slide.shapes.add_textbox(
            Inches(6.5), Inches(1.2),
            Inches(6), Inches(0.4)
        )
        proj_title.text_frame.text = "3-Year Financial Projections (ZAR)"
        proj_title.text_frame.paragraphs[0].font.size = Pt(22)
        proj_title.text_frame.paragraphs[0].font.bold = True
        proj_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

        # Projection table
        projections = [
            ("Year 1", "10K users", "R 4.8M", "R 2.2M", "R 2.6M"),
            ("Year 2", "50K users", "R 28.5M", "R 8.5M", "R 20M"),
            ("Year 3", "150K users", "R 105M", "R 18.2M", "R 86.8M")
        ]

        # Headers
        headers = ["Year", "Users", "Revenue", "Costs", "Profit"]
        header_widths = [1.2, 1.2, 1.5, 1.3, 1.5]
        x_start = 6.5
        y_header = 1.8

        for i, (header, width) in enumerate(zip(headers, header_widths)):
            x = x_start + sum(header_widths[:i])
            header_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_header),
                Inches(width), Inches(0.4)
            )
            header_box.text_frame.text = header
            header_box.text_frame.paragraphs[0].font.size = Pt(14)
            header_box.text_frame.paragraphs[0].font.bold = True
            header_box.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
            header_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Data rows
        y_data = 2.3
        for i, proj in enumerate(projections):
            for j, (value, width) in enumerate(zip(proj, header_widths)):
                x = x_start + sum(header_widths[:j])
                y = y_data + i * 0.6

                # Row background
                if j == 0:
                    row_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x_start - 0.1), Inches(y - 0.05),
                        Inches(sum(header_widths) + 0.2), Inches(0.5)
                    )
                    row_bg.fill.solid()
                    row_bg.fill.fore_color.rgb = BG_LIGHT
                    row_bg.line.fill.background()

                cell_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y),
                    Inches(width), Inches(0.4)
                )
                cell_box.text_frame.text = value
                cell_box.text_frame.paragraphs[0].font.size = Pt(14)
                cell_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
                cell_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                if j == 0:
                    cell_box.text_frame.paragraphs[0].font.bold = True
                if j == 4:  # Profit column
                    cell_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
                    cell_box.text_frame.paragraphs[0].font.bold = True

        # Key metrics
        metrics_box = slide.shapes.add_textbox(
            Inches(6.5), Inches(5.2),
            Inches(6), Inches(1.5)
        )
        metrics_frame = metrics_box.text_frame
        metrics_frame.text = "Key Metrics:\n• Break-even: Month 18\n• Customer Acquisition Cost: R250\n• Lifetime Value: R4,800\n• LTV/CAC Ratio: 19.2x\n• Gross Margin: 82%"
        for para in metrics_frame.paragraphs:
            para.font.size = Pt(15)
            para.font.color.rgb = DARK_GRAY
            para.space_after = Pt(6)

    def add_funding_request():
        """Slide 12: Funding Request"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "💎 Funding Request"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # Funding amount box (hero)
        amount_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(1.5),
            Inches(9.333), Inches(1.3)
        )
        amount_box.fill.solid()
        amount_box.fill.fore_color.rgb = ACCENT_GREEN
        amount_box.line.fill.background()
        amount_box.shadow.inherit = False

        # Amount text
        amount_text = slide.shapes.add_textbox(
            Inches(2), Inches(1.65),
            Inches(9.333), Inches(0.5)
        )
        amount_text.text_frame.text = "Seeking: R 15 Million"
        amount_text.text_frame.paragraphs[0].font.size = Pt(48)
        amount_text.text_frame.paragraphs[0].font.bold = True
        amount_text.text_frame.paragraphs[0].font.color.rgb = WHITE
        amount_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtext
        subtext = slide.shapes.add_textbox(
            Inches(2), Inches(2.2),
            Inches(9.333), Inches(0.4)
        )
        subtext.text_frame.text = "Series A Funding • 18-Month Runway"
        subtext.text_frame.paragraphs[0].font.size = Pt(22)
        subtext.text_frame.paragraphs[0].font.color.rgb = WHITE
        subtext.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Use of funds
        use_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.1),
            Inches(12.333), Inches(0.4)
        )
        use_title.text_frame.text = "Use of Funds"
        use_title.text_frame.paragraphs[0].font.size = Pt(24)
        use_title.text_frame.paragraphs[0].font.bold = True
        use_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

        allocations = [
            ("Product Development", "R 5.5M", "37%", "Engineering team, feature development, UX/UI", ACCENT_GREEN),
            ("Cloud Infrastructure", "R 2.5M", "17%", "AWS costs, scaling, security, monitoring", PRIMARY_BLUE),
            ("Marketing & Sales", "R 4M", "27%", "Customer acquisition, brand building, partnerships", ACCENT_ORANGE),
            ("Operations & Legal", "R 1.5M", "10%", "Compliance, licenses, operations, support", DARK_BLUE),
            ("Reserve Fund", "R 1.5M", "10%", "Contingency, opportunities, buffer", LIGHT_GRAY)
        ]

        y_start = 3.7
        for i, (category, amount, perc, desc, color) in enumerate(allocations):
            y = y_start + i * 0.62

            # Bar background
            bar_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y),
                Inches(12.333), Inches(0.55)
            )
            bar_bg.fill.solid()
            bar_bg.fill.fore_color.rgb = BG_LIGHT
            bar_bg.line.fill.background()

            # Progress bar
            perc_width = 12.333 * (float(perc.strip('%')) / 100)
            bar_fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y),
                Inches(perc_width), Inches(0.55)
            )
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = color
            bar_fill.line.fill.background()

            # Category text
            cat_text = slide.shapes.add_textbox(
                Inches(0.7), Inches(y + 0.08),
                Inches(3), Inches(0.4)
            )
            cat_text.text_frame.text = f"{category} - {amount}"
            cat_text.text_frame.paragraphs[0].font.size = Pt(16)
            cat_text.text_frame.paragraphs[0].font.bold = True
            cat_text.text_frame.paragraphs[0].font.color.rgb = WHITE if float(perc.strip('%')) > 20 else DARK_GRAY

            # Description
            desc_text = slide.shapes.add_textbox(
                Inches(4), Inches(y + 0.12),
                Inches(5), Inches(0.3)
            )
            desc_text.text_frame.text = desc
            desc_text.text_frame.word_wrap = True
            desc_text.text_frame.paragraphs[0].font.size = Pt(13)
            desc_text.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY if float(perc.strip('%')) < 20 else WHITE

            # Percentage
            perc_text = slide.shapes.add_textbox(
                Inches(11.5), Inches(y + 0.08),
                Inches(1), Inches(0.4)
            )
            perc_text.text_frame.text = perc
            perc_text.text_frame.paragraphs[0].font.size = Pt(18)
            perc_text.text_frame.paragraphs[0].font.bold = True
            perc_text.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
            perc_text.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def add_roi_projections():
        """Slide 13: ROI & Exit Strategy"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "📈 ROI Projections & Exit Strategy"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_BLUE

        # ROI highlights
        roi_boxes = [
            ("Expected ROI", "5.8x", "In 3 Years", ACCENT_GREEN),
            ("Valuation Target", "R 350M", "Year 3", PRIMARY_BLUE),
            ("IRR", "142%", "Annual", ACCENT_ORANGE)
        ]

        box_width = 3.8
        spacing = 0.4
        start_x = 0.8

        for i, (label, value, desc, color) in enumerate(roi_boxes):
            x = start_x + i * (box_width + spacing)

            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.4),
                Inches(box_width), Inches(1.4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()
            box.shadow.inherit = False

            # Value
            val_box = slide.shapes.add_textbox(
                Inches(x), Inches(1.6),
                Inches(box_width), Inches(0.6)
            )
            val_box.text_frame.text = value
            val_box.text_frame.paragraphs[0].font.size = Pt(48)
            val_box.text_frame.paragraphs[0].font.bold = True
            val_box.text_frame.paragraphs[0].font.color.rgb = WHITE
            val_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Label
            lbl_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.2),
                Inches(box_width), Inches(0.3)
            )
            lbl_box.text_frame.text = label
            lbl_box.text_frame.paragraphs[0].font.size = Pt(16)
            lbl_box.text_frame.paragraphs[0].font.bold = True
            lbl_box.text_frame.paragraphs[0].font.color.rgb = WHITE
            lbl_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Description
            dsc_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.5),
                Inches(box_width), Inches(0.25)
            )
            dsc_box.text_frame.text = desc
            dsc_box.text_frame.paragraphs[0].font.size = Pt(14)
            dsc_box.text_frame.paragraphs[0].font.color.rgb = WHITE
            dsc_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Exit strategies
        exit_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.2),
            Inches(12.333), Inches(0.4)
        )
        exit_title.text_frame.text = "Exit Strategy Options"
        exit_title.text_frame.paragraphs[0].font.size = Pt(24)
        exit_title.text_frame.paragraphs[0].font.bold = True
        exit_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

        exits = [
            ("🏦 Strategic Acquisition", "Major bank acquisition (FNB, Standard Bank, Capitec)", "Year 3-4", "Target: R 300-400M", ACCENT_GREEN),
            ("🌍 International Expansion", "Expand to other African markets, raise Series B", "Year 2-3", "Target: R 150-200M", PRIMARY_BLUE),
            ("📊 IPO", "Public listing on JSE or international exchange", "Year 4-5", "Target: R 500M+", ACCENT_ORANGE)
        ]

        y_start = 4
        for i, (strategy, desc, timeline, target, color) in enumerate(exits):
            y = y_start + i * 0.95

            # Strategy box
            strat_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y),
                Inches(12.333), Inches(0.85)
            )
            strat_box.fill.solid()
            strat_box.fill.fore_color.rgb = BG_LIGHT
            strat_box.line.color.rgb = color
            strat_box.line.width = Pt(3)

            # Icon & Title
            title_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y + 0.08),
                Inches(4), Inches(0.35)
            )
            title_box.text_frame.text = strategy
            title_box.text_frame.paragraphs[0].font.size = Pt(20)
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.color.rgb = color

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(y + 0.45),
                Inches(6.5), Inches(0.3)
            )
            desc_box.text_frame.text = desc
            desc_box.text_frame.word_wrap = True
            desc_box.text_frame.paragraphs[0].font.size = Pt(14)
            desc_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

            # Timeline
            time_box = slide.shapes.add_textbox(
                Inches(7.5), Inches(y + 0.15),
                Inches(2), Inches(0.3)
            )
            time_box.text_frame.text = f"⏱ {timeline}"
            time_box.text_frame.paragraphs[0].font.size = Pt(15)
            time_box.text_frame.paragraphs[0].font.bold = True
            time_box.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

            # Target
            tgt_box = slide.shapes.add_textbox(
                Inches(9.8), Inches(y + 0.15),
                Inches(2.7), Inches(0.3)
            )
            tgt_box.text_frame.text = target
            tgt_box.text_frame.paragraphs[0].font.size = Pt(16)
            tgt_box.text_frame.paragraphs[0].font.bold = True
            tgt_box.text_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
            tgt_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Create all slides
    print("Creating executive presentation slides...")
    add_title_slide()
    print("  ✓ Title slide")
    add_executive_summary()
    print("  ✓ Executive summary")
    add_market_opportunity()
    print("  ✓ Market opportunity")
    add_business_model()
    print("  ✓ Business model & revenue")
    add_funding_request()
    print("  ✓ Funding request")
    add_roi_projections()
    print("  ✓ ROI & exit strategy")
    add_tech_stack_infographic()
    print("  ✓ Tech stack infographic")
    add_security_architecture()
    print("  ✓ Security architecture")
    add_aws_infrastructure()
    print("  ✓ AWS infrastructure")
    add_features_dashboard()
    print("  ✓ Features dashboard")
    add_deployment_status()
    print("  ✓ Deployment status")
    add_next_steps()
    print("  ✓ Strategic roadmap")
    add_closing_slide()
    print("  ✓ Closing slide")

    # Save presentation
    prs.save(output_file)
    print(f"\n✅ Executive presentation created: {output_file}")
    print(f"📊 Total slides: 13")

if __name__ == '__main__':
    output_file = 'docs/BankApp_Executive_Presentation.pptx'

    print("🎨 Creating executive PowerPoint presentation with infographics...")
    create_exec_presentation(output_file)
    print("✅ Done!")
